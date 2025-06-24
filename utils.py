import os
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.vectorstores import FAISS
from ctransformers import AutoModelForCausalLM
from transformers import AutoTokenizer, pipeline
import os

# --- New Imports for File Parsing ---
from docx import Document # For .docx files
import fitz # PyMuPDF for .pdf files
import io
from PIL import Image # For image handling in OCR
import pytesseract # For OCR
import pandas as pd # For .xls, .xlsx, .csv files
from pptx import Presentation # For .pptx files




import configs

# Function to get all supported files in the projects folder and its subfolders
def get_files(path):
    """
    Retrieves all supported files in the specified folder and its subfolders

    Args:
        path (str): Path to the folder containing the files

    Returns:
        list: List of paths to supported files
    """
     
    try:
        supported_formats = [".pdf", ".docx", ".ipynb", ".py", ".md", ".pptx", ".xls", ".xlsx", ".csv"]
        supported_files = []
        for root, dirs, files in os.walk(path):
            for file in files:
                if file.endswith(tuple(supported_formats)):
                    supported_files.append(os.path.join(root, file))
        return supported_files
    except Exception as e:
        print(e)
        return []

# Function to parse multiple files
def parse_files(files):
    """
    Parses multiple files and returns concatenated text

    Args:
        files (list): List of file paths to be parsed

    Returns:
        tuple: Concatenated text and number of files successfully parsed
    """
    texts = ""
    i = 0
    for file in files:
        try:
            parsed_text = parse_file(file)
            print("%s parsed successfully" % file)
            texts += "\n" + parsed_text
            i += 1
        except Exception as e:
            print("%s failed to parse" % file)
            print(f"Error: {e}")
    return texts, i

# Function to create the vector database with Sentence Transformers
def create_vector_db(texts):
    """
    Creates a vector database with Sentence Transformers for text embeddings.

    Args:
        texts (str): Text data to be processed.

    Returns:
        None
    """
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=20)
    split_texts = text_splitter.split_text(texts)

    embeddings = HuggingFaceEmbeddings(model_name='sentence-transformers/all-MiniLM-L6-v2',
                                       model_kwargs={'device': 'cpu'})
    db = FAISS.from_texts(split_texts, embeddings)
    db.save_local(configs.DB_FAISS_PATH)

# Function to load the model and tokenizer 
def load_model_and_tokenizer(model_id):
    """
    Loads the language model and tokenizer for the pipeline using ctransformers.

    Args:
        model_id (str): Hugging Face ID of the GGUF model (e.g., "TheBloke/Llama-2-7B-Chat-GGUF").

    Returns:
        pipeline: Initialized text-generation pipeline.
    """
    model = AutoModelForCausalLM.from_pretrained(
        model_id,
        model_file="llama-2-7b-chat.Q4_K_M.gguf", # Verify this filename on Hugging Face
        model_type="llama",
        gpu_layers=0,
        hf=True
    )
    
    return model


# --- File Parsing Functions  ---

def parse_pdf_searchable(file_path):
    """Extracts all searchable text from a PDF file using PyMuPDF."""
    text_content = ""
    try:
        document = fitz.open(file_path)
        for page_num in range(document.page_count):
            page = document.load_page(page_num)
            text_content += page.get_text()
        document.close()
    except Exception as e:
        raise Exception(f"Failed to parse searchable PDF: {e}") from e
    return text_content

def parse_pdf_scanned_ocr(file_path):
    """Extracts text from a scanned PDF using PyMuPDF to extract images and Pytesseract for OCR."""
    text_content = ""
    try:
        document = fitz.open(file_path)
        for page_num in range(document.page_count):
            page = document.load_page(page_num)
            image_list = page.get_images(full=True)

            if not image_list:
                pix = page.get_pixmap()
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                text_content += pytesseract.image_to_string(img, lang='eng') + "\n"
            else:
                for img_index, img_info in enumerate(image_list):
                    xref = img_info[0]
                    base_image = document.extract_image(xref)
                    image_bytes = base_image["image"]

                    image = Image.open(io.BytesIO(image_bytes))
                    text_from_image = pytesseract.image_to_string(image, lang='eng')
                    text_content += text_from_image + "\n"
        document.close()
    except Exception as e:
        raise Exception(f"Failed to parse scanned PDF with OCR: {e}") from e
    return text_content

def parse_docx(file_path):
    """Extracts text from a .docx file using python-docx."""
    try:
        document = Document(file_path)
        full_text = []
        for paragraph in document.paragraphs:
            full_text.append(paragraph.text)
        return '\n'.join(full_text)
    except Exception as e:
        raise Exception(f"Failed to parse DOCX file: {e}") from e

def parse_excel(file_path):
    """Extracts text from .xls and .xlsx files using pandas."""
    try:
        df = pd.read_excel(file_path)
        return df.to_string(index=False)
    except Exception as e:
        raise Exception(f"Failed to parse Excel file: {e}") from e

def parse_csv(file_path):
    """Extracts text from .csv files using pandas."""
    try:
        df = pd.read_csv(file_path)
        return df.to_string(index=False)
    except Exception as e:
        raise Exception(f"Failed to parse CSV file: {e}") from e

def parse_pptx(file_path):
    """Extracts text from .pptx files."""
    text_content = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        return '\n'.join(text_content)
    except Exception as e:
        raise Exception(f"Failed to parse PPTX file: {e}") from e

def parse_file(file_name):
    """
    Parses a file based on its extension using different libraries.
    """
    file_extension = os.path.splitext(file_name)[1].lower()

    if file_extension == ".pdf":
        text = parse_pdf_searchable(file_name)
        if len(text.strip()) < 50:
            print(f"น้อยข้อความที่พบใน {file_name} (searchable). กำลังลอง OCR...")
            text = parse_pdf_scanned_ocr(file_name)
        return text
    elif file_extension == ".docx":
        return parse_docx(file_name)
    elif file_extension == ".pptx":
        return parse_pptx(file_name)
    elif file_extension in [".xls", ".xlsx"]:
        return parse_excel(file_name)
    elif file_extension == ".csv":
        return parse_csv(file_name)
    elif file_extension in [".ipynb", ".py", ".md"]:
        try:
            with open(file_name, "r", encoding='utf-8') as f:
                data = f.read()
            return data
        except Exception as e:
            raise Exception(f"Failed to parse file using python standard library: {e}") from e
    else:
        raise Exception(f"File type not supported: {file_extension}")


def streamlit_parse_file(uploaded_file):
    """
    Parses a file uploaded in a Streamlit application.
    """
    file_extension = uploaded_file.name.split(".")[-1].lower()
    
    original_pos = uploaded_file.tell()
    uploaded_file.seek(0)

    import tempfile
    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_extension}") as temp_file:
        temp_file.write(uploaded_file.read())
        temp_file_path = temp_file.name

    uploaded_file.seek(original_pos)

    text_content = None
    try:
        if file_extension == "pdf":
            text_content = parse_pdf_searchable(temp_file_path)
            if len(text_content.strip()) < 50:
                print(f"Little searchable text found in uploaded PDF. Attempting OCR...")
                text_content = parse_pdf_scanned_ocr(temp_file_path)
        elif file_extension == "docx":
            text_content = parse_docx(temp_file_path)
        elif file_extension == "pptx":
            text_content = parse_pptx(temp_file_path)
        elif file_extension in ["xls", "xlsx"]:
            text_content = parse_excel(temp_file_path)
        elif file_extension == "csv":
            text_content = parse_csv(temp_file_path)
        elif file_extension in ["ipynb", "py", "md"]:
            uploaded_file.seek(0)
            text_content = uploaded_file.read().decode("utf-8")
        else:
            raise Exception("File type not supported")
    except Exception as e:
        raise Exception(f"Failed to parse uploaded file: {str(e)}")
    finally:
        if os.path.exists(temp_file_path):
            os.remove(temp_file_path)
            
    return text_content